import streamlit as st
from PyPDF2 import PdfReader
import pptx
from openai import OpenAI
# # import os
# from dotenv import find_dotenv, load_dotenv

# Import the OpenAI API key

API_KEY = st.secrets.OPENAI_API_KEY

# Initialize OpenAI client
client = OpenAI(api_key=API_KEY)


# Set page title and configure layout
st.set_page_config( layout="wide")
# Custom CSS for green shades
st.markdown(
    """
    <style>
    .stApp {
        background-color: #e6f5e6;
    }
    .stButton>button {
        background-color: #4CAF50;
        color: white;
        border-radius: 5px;
        border: none;
        padding: 10px 24px;
        font-size: 16px;
    }
    .stButton>button:hover {
        background-color: #45a049;
    }
    .stFileUploader>div>div>div>button {
        background-color: #4CAF50;
        color: white;
    }
    .stFileUploader>div>div>div>button:hover {
        background-color: #45a049;
    }
    </style>
    """,
    unsafe_allow_html=True
)

# Title of the webpage
st.title("Technical Evaluation Validator")

# File uploader
uploaded_file = st.file_uploader("Upload a PDF or PowerPoint file", type=["pdf", "pptx"])

# Function to extract text from PDF
def extract_text_from_pdf(file):
    pdf_reader = PdfReader(file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

# Function to extract text from PowerPoint
def extract_text_from_pptx(file):
    presentation = pptx.Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to summarize text using OpenAI
def summarize_text(text):
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are a helpful assistant that summarizes text."},
            {"role": "user", "content": f"Summarize the following text:\n{text}\n\nSummary:"}
        ],
        max_tokens=150,
        temperature=0.5
    )
    return response.choices[0].message.content.strip()

# Analyze button
if st.button("Analyze"):
    st.markdown(
        """
        <style>
        hr {
            border: 1px solid #ffffff; /* Change the color to green */
        }
        </style>
        """,
        unsafe_allow_html=True
    )
    st.divider()
    if uploaded_file is not None:
        # file_details = {"filename": uploaded_file.name, "filetype": uploaded_file.type, "filesize": uploaded_file.size}
        # st.write(file_details)

        # Extract text based on file type
        if uploaded_file.type == "application/pdf":
            text = extract_text_from_pdf(uploaded_file)
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            text = extract_text_from_pptx(uploaded_file)
        else:
            st.error("Unsupported file type")
            text = ""

        if text:
            # Create two columns for layout
            col1, col2 = st.columns(2)

            # Display extracted text in the left column
            with col1:
                st.subheader("Extracted Text")
                st.write(text)

            # Summarize the text and display the summary in the right column
            with col2:
                st.subheader("Validation Report")
                with st.spinner("Validating..."):
                    summary = summarize_text(text[:1000])  # Summarize the first 1000 characters
                    st.write(summary)
    else:
        st.warning("Please upload a file first.")
